"""Convert Marp-flavored PITCH.md to a .pptx deck.

Splits on `---` slide boundaries (after frontmatter), parses each slide as
title + body blocks (paragraphs, code, tables, lists, blockquotes), and
renders each block to a real PowerPoint shape.

Run:
    python3 _md_to_pptx.py PITCH.md PITCH.pptx
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# ── colors ─────────────────────────────────────────────────────────────
INK   = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x0B, 0x5F, 0xFF)
CODE_BG = RGBColor(0xF4, 0xF4, 0xF4)
CODE_FG = RGBColor(0x1A, 0x1A, 0x1A)
HEADER_BG = RGBColor(0xE8, 0xEE, 0xFB)

# ── slide geometry (16:9) ───────────────────────────────────────────────
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN_X = Inches(0.6)
TITLE_Y  = Inches(0.35)
TITLE_H  = Inches(0.8)
BODY_Y   = Inches(1.25)
BODY_H   = SLIDE_H - BODY_Y - Inches(0.4)
BODY_W   = SLIDE_W - 2 * MARGIN_X


# ── markdown parsing ────────────────────────────────────────────────────

def strip_frontmatter(text: str) -> str:
    if text.startswith("---"):
        m = re.match(r"^---\s*\n.*?\n---\s*\n", text, re.DOTALL)
        if m:
            return text[m.end():]
    return text


def split_slides(text: str) -> list[str]:
    """Split on lines that are exactly '---' (slide boundary in Marp)."""
    parts = re.split(r"^---\s*$", text, flags=re.MULTILINE)
    return [p.strip("\n") for p in parts if p.strip()]


# Block types: ('title', text), ('para', text), ('code', text),
# ('table', rows), ('list', items), ('quote', text), ('hr', None)
def parse_blocks(slide_md: str):
    lines = slide_md.split("\n")
    blocks = []
    i = 0
    title = None

    while i < len(lines):
        line = lines[i]

        # Title — first # or ## (or ###)
        if title is None and line.startswith("#"):
            m = re.match(r"^(#{1,3})\s+(.*)$", line)
            if m:
                title = m.group(2).strip()
                i += 1
                continue

        # Code block
        if line.startswith("```"):
            j = i + 1
            buf = []
            while j < len(lines) and not lines[j].startswith("```"):
                buf.append(lines[j])
                j += 1
            blocks.append(("code", "\n".join(buf)))
            i = j + 1
            continue

        # Table — line starts with | and next line is separator
        if line.startswith("|") and i + 1 < len(lines) and re.match(r"^\|[\s:|-]+\|?\s*$", lines[i+1]):
            rows = []
            # header
            rows.append(_parse_table_row(line))
            j = i + 2  # skip separator
            while j < len(lines) and lines[j].startswith("|"):
                rows.append(_parse_table_row(lines[j]))
                j += 1
            blocks.append(("table", rows))
            i = j
            continue

        # Blockquote
        if line.startswith("> "):
            buf = []
            while i < len(lines) and lines[i].startswith(">"):
                buf.append(lines[i].lstrip("> ").rstrip())
                i += 1
            blocks.append(("quote", " ".join(buf)))
            continue

        # List (- or 1. or - [ ])
        if re.match(r"^\s*([-*]|\d+\.)\s+", line):
            items = []
            while i < len(lines) and (re.match(r"^\s*([-*]|\d+\.)\s+", lines[i]) or lines[i].startswith("  ")):
                m = re.match(r"^(\s*)([-*]|\d+\.)\s+(.*)$", lines[i])
                if m:
                    indent = len(m.group(1)) // 2
                    items.append((indent, m.group(3).rstrip()))
                elif lines[i].strip() == "":
                    break
                else:
                    if items:
                        items[-1] = (items[-1][0], items[-1][1] + " " + lines[i].strip())
                i += 1
            blocks.append(("list", items))
            continue

        # Blank line
        if line.strip() == "":
            i += 1
            continue

        # Paragraph — gather until blank line / list / table / code / heading
        buf = [line]
        i += 1
        while i < len(lines):
            nxt = lines[i]
            if (nxt.strip() == "" or nxt.startswith("#") or nxt.startswith("```")
                    or nxt.startswith("|") or nxt.startswith("> ")
                    or re.match(r"^\s*([-*]|\d+\.)\s+", nxt)):
                break
            buf.append(nxt)
            i += 1
        blocks.append(("para", " ".join(b.strip() for b in buf)))

    return title, blocks


def _parse_table_row(line: str) -> list[str]:
    # | a | b | c |   → ['a', 'b', 'c']
    parts = line.strip().strip("|").split("|")
    return [p.strip() for p in parts]


# ── inline markdown → runs ──────────────────────────────────────────────

INLINE_RE = re.compile(
    r"(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`|\[[^\]]+\]\([^)]+\))"
)


def add_inline(p, text: str, *, base_size: int = 18, base_color=INK, base_bold=False):
    """Append runs to paragraph p, parsing **bold**, *italic*, `code`, [link](url)."""
    if not text:
        return
    pos = 0
    for m in INLINE_RE.finditer(text):
        if m.start() > pos:
            _run(p, text[pos:m.start()], base_size, base_color, bold=base_bold)
        token = m.group(1)
        if token.startswith("**"):
            _run(p, token[2:-2], base_size, base_color, bold=True)
        elif token.startswith("*"):
            _run(p, token[1:-1], base_size, base_color, italic=True, bold=base_bold)
        elif token.startswith("`"):
            _run(p, token[1:-1], base_size, CODE_FG, mono=True, bold=base_bold)
        elif token.startswith("["):
            lm = re.match(r"\[([^\]]+)\]\(([^)]+)\)", token)
            if lm:
                _run(p, lm.group(1), base_size, ACCENT, bold=base_bold, underline=True)
        pos = m.end()
    if pos < len(text):
        _run(p, text[pos:], base_size, base_color, bold=base_bold)


def _run(p, text, size, color, *, bold=False, italic=False, mono=False, underline=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.underline = underline
    if mono:
        r.font.name = "Consolas"


# ── rendering ───────────────────────────────────────────────────────────

def render(prs, slides):
    blank_layout = prs.slide_layouts[6]  # truly blank
    for idx, slide_md in enumerate(slides):
        title, blocks = parse_blocks(slide_md)
        slide = prs.slides.add_slide(blank_layout)

        # accent stripe
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H)
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        # title
        if title:
            tb = slide.shapes.add_textbox(MARGIN_X, TITLE_Y, BODY_W, TITLE_H)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            add_inline(p, title, base_size=30, base_bold=True, base_color=INK)

            # underline
            ln = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, MARGIN_X, TITLE_Y + Inches(0.85),
                Inches(0.7), Inches(0.04),
            )
            ln.fill.solid()
            ln.fill.fore_color.rgb = ACCENT
            ln.line.fill.background()

        # body
        cur_y = BODY_Y
        for kind, payload in blocks:
            cur_y = render_block(slide, kind, payload, cur_y, BODY_W, MARGIN_X)
            if cur_y >= SLIDE_H - Inches(0.3):
                break  # overflow — clip rather than spill

        # page number
        pn = slide.shapes.add_textbox(SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.4),
                                      Inches(0.8), Inches(0.3))
        pn.text_frame.margin_left = pn.text_frame.margin_right = 0
        p = pn.text_frame.paragraphs[0]
        p.alignment = 2  # right
        _run(p, f"{idx+1}", 9, MUTED)


def render_block(slide, kind, payload, y, width, x):
    if kind == "para":
        return _render_para(slide, payload, y, width, x)
    if kind == "code":
        return _render_code(slide, payload, y, width, x)
    if kind == "list":
        return _render_list(slide, payload, y, width, x)
    if kind == "table":
        return _render_table(slide, payload, y, width, x)
    if kind == "quote":
        return _render_quote(slide, payload, y, width, x)
    return y


def _emu_to_inches(emu):
    return emu / 914400.0


def _render_para(slide, text, y, width, x):
    # estimate height by line count at ~11 chars/inch
    cpl = max(1, int(_emu_to_inches(width) * 11))
    lines = max(1, (len(text) + cpl - 1) // cpl)
    h = Inches(0.36 * lines + 0.1)
    tb = slide.shapes.add_textbox(x, y, width, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    add_inline(p, text, base_size=18, base_color=INK)
    return y + h + Inches(0.1)


def _render_code(slide, text, y, width, x):
    lines = text.split("\n")
    n = len(lines)
    line_h = 0.22
    h = Inches(line_h * n + 0.25)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, h)
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    box.line.width = Pt(0.5)

    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  width - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(12)
        r.font.color.rgb = CODE_FG
    return y + h + Inches(0.12)


def _render_list(slide, items, y, width, x):
    n = len(items)
    h = Inches(0.34 * n + 0.1)
    tb = slide.shapes.add_textbox(x, y, width, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, (indent, item) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        p.space_before = Pt(2)
        # bullet
        bullet = "  " * indent + ("• " if indent == 0 else "– ")
        text = item
        # checkboxes: "[ ] foo" / "[x] foo"
        m = re.match(r"^\[([ xX])\]\s+(.*)$", text)
        if m:
            text = m.group(2)
            bullet = "  " * indent + ("☑ " if m.group(1).lower() == "x" else "☐ ")
        _run(p, bullet, 16, MUTED)
        add_inline(p, text, base_size=16, base_color=INK)
    return y + h + Inches(0.05)


def _render_table(slide, rows, y, width, x):
    if not rows:
        return y
    nrows = len(rows)
    ncols = len(rows[0])
    row_h = Inches(0.42)
    h = row_h * nrows
    if h > Inches(5.5):
        # shrink rows if huge
        row_h = Inches(5.5 / nrows)
        h = row_h * nrows
    tbl_shape = slide.shapes.add_table(nrows, ncols, x, y, width, h)
    tbl = tbl_shape.table

    # roughly proportional column widths — first col gets a bit more
    if ncols == 2:
        tbl.columns[0].width = int(width * 0.32)
        tbl.columns[1].width = width - tbl.columns[0].width
    elif ncols == 3:
        tbl.columns[0].width = int(width * 0.10)
        tbl.columns[1].width = int(width * 0.45)
        tbl.columns[2].width = width - tbl.columns[0].width - tbl.columns[1].width

    for r_idx, row in enumerate(rows):
        for c_idx, cell_md in enumerate(row[:ncols]):
            cell = tbl.cell(r_idx, c_idx)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = HEADER_BG if r_idx == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = ""  # clear default
            p = tf.paragraphs[0]
            add_inline(p, cell_md, base_size=13,
                       base_color=INK, base_bold=(r_idx == 0))
    return y + h + Inches(0.15)


def _render_quote(slide, text, y, width, x):
    h = Inches(0.7)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(x + Inches(0.25), y, width - Inches(0.25), h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    add_inline(p, text, base_size=18, base_color=MUTED)
    p.runs[0].font.italic = True
    return y + h + Inches(0.1)


# ── main ────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) != 3:
        print("usage: _md_to_pptx.py <input.md> <output.pptx>", file=sys.stderr)
        sys.exit(1)
    src = Path(sys.argv[1])
    dst = Path(sys.argv[2])

    text = strip_frontmatter(src.read_text(encoding="utf-8"))
    slides = split_slides(text)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    render(prs, slides)
    prs.save(dst)
    print(f"wrote {dst} — {len(slides)} slides")


if __name__ == "__main__":
    main()
